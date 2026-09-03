from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUTS = [
    ROOT / "documentation" / "Solution_Documentation.pptx",
    ROOT / "submission_package" / "documentation" / "Solution_Documentation.pptx",
]

NAVY = RGBColor(20, 54, 88)
BLUE = RGBColor(35, 94, 142)
CYAN = RGBColor(33, 158, 188)
ORANGE = RGBColor(238, 135, 52)
INK = RGBColor(32, 42, 52)
MUTED = RGBColor(91, 105, 117)
PALE = RGBColor(240, 246, 249)
WHITE = RGBColor(255, 255, 255)
FONT = "Aptos"

SLIDES = [
    (
        "01", "Approach", "Compact recognition designed for embedded inference",
        [
            ("Algorithm / model", "A 5-block depthwise-separable CNN (16->24->32->48->64) with learned 6x6 per-channel spatial pooling and a Dense(17) classifier."),
            ("Why this model", "Depthwise convolutions provide strong image features with few parameters; learned spatial pooling preserves useful sign location cues while keeping flash and RAM usage low."),
        ],
    ),
    (
        "02", "Training", "Orientation-aware learning on the supplied dataset",
        [
            ("How it was trained", "48x48 RGB images; class-aware oversampling; +/-6 degree rotation, translation, zoom, brightness and contrast augmentation; class weights; Adam with cosine learning-rate decay; label smoothing 0.05 and light MixUp (alpha 0.10)."),
            ("Training / validation result", "Quantization-aware training (25 epochs) evaluated every checkpoint as INT8. The selected checkpoint achieved 352/364 correct on validation: 96.70%."),
        ],
    ),
    (
        "03", "Optimization", "Full-integer deployment without measured accuracy loss",
        [
            ("Technique", "Quantization-aware training followed by full-integer INT8 TFLite conversion with INT8 input and output. The best checkpoint was selected using exact INT8 validation accuracy."),
            ("Model size", "Before: 234,708-byte Keras artifact. After: 30,168-byte TFLite model (87.1% smaller)."),
            ("Accuracy impact", "Selected checkpoint before conversion: 96.70%. INT8 model after conversion: 96.70% (352/364). No measured validation loss."),
        ],
    ),
    (
        "04", "Embedded Deployment", "TensorFlow Lite Micro integrated with Zephyr RTOS",
        [
            ("Zephyr integration", "The TFLite model is embedded as a C byte array. A static TFLM interpreter and 40 KB tensor arena register only the required operators. Firmware receives a 6,912-byte RGB frame over UART and returns the official class ID."),
            ("Running in QEMU", "Build the Zephyr qemu_x86 target, boot it with QEMU, then use qemu_uart_test.py to send the 0xAA 0x55 marker plus image bytes and read the predicted ID. Predictions were checked against desktop TFLite output."),
        ],
    ),
    (
        "05", "Performance on Validation Set in QEMU", "Measured with the deployed v16 INT8 firmware",
        [
            ("Accuracy", "352/364 correct (96.70%), with exact QEMU prediction parity on the tested validation images."),
            ("Model size", "30,168 bytes (29.5 KB), 11,417 parameters."),
            ("RAM usage", "TFLM tensor arena: 28,436 / 40,960 bytes used. Zephyr build RAM: 279,164 bytes."),
            ("Inference time", "Approximately 82-106 ms per frame in qemu_x86. QEMU timing is emulated and is not physical ESP32-S3 latency."),
        ],
    ),
]


def add_textbox(slide, left, top, width, height, text, size, color, bold=False,
                font=FONT, align=PP_ALIGN.LEFT, margin=0.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_slide(prs, number, title, subtitle, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE

    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), prs.slide_height)
    rail.fill.solid()
    rail.fill.fore_color.rgb = CYAN
    rail.line.fill.background()

    add_textbox(slide, Inches(0.68), Inches(0.42), Inches(1.05), Inches(0.4),
                number, 14, ORANGE, True)
    add_textbox(slide, Inches(0.68), Inches(0.78), Inches(11.8), Inches(0.72),
                title, 30, NAVY, True)
    add_textbox(slide, Inches(0.7), Inches(1.48), Inches(11.5), Inches(0.42),
                subtitle, 15, MUTED)

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.02),
                                  Inches(11.9), Inches(0.025))
    rule.fill.solid()
    rule.fill.fore_color.rgb = CYAN
    rule.line.fill.background()

    count = len(items)
    gap = 0.16
    available = 4.75
    card_height = (available - gap * (count - 1)) / count
    y = 2.25
    for index, (label, body) in enumerate(items):
        accent = ORANGE if index % 2 == 0 else CYAN
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.7), Inches(y), Inches(11.9), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = PALE
        card.line.color.rgb = RGBColor(218, 229, 235)
        card.line.width = Pt(0.8)
        card.adjustments[0] = 0.08

        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.7), Inches(y), Inches(0.09), Inches(card_height))
        marker.fill.solid()
        marker.fill.fore_color.rgb = accent
        marker.line.fill.background()

        add_textbox(slide, Inches(1.02), Inches(y + 0.12), Inches(2.45),
                    Inches(card_height - 0.24), label, 14, BLUE, True)
        body_size = 15 if count <= 3 else 13.5
        add_textbox(slide, Inches(3.42), Inches(y + 0.1), Inches(8.82),
                    Inches(card_height - 0.2), body, body_size, INK)
        y += card_height + gap

    add_textbox(slide, Inches(10.85), Inches(7.1), Inches(1.4), Inches(0.22),
                "AI ARENA 2026", 9, MUTED, True, align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for slide_data in SLIDES:
        add_slide(prs, *slide_data)
    for output in OUTPUTS:
        output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output)
        print(f"Wrote {output}")


if __name__ == "__main__":
    build()
